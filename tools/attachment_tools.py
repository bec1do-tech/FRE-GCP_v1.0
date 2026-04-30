"""
FRE GCP v1.0 — File Attachment Tools
======================================
Lets users drag-and-drop or attach files directly in the ADK chat UI.

Two tools are provided:

1. extract_office_document_text(file_base64, filename)
   ─────────────────────────────────────────────────
   Gemini can natively read PDF and image attachments (they arrive as
   inline_data in the message).  For binary Office formats (DOCX, PPTX,
   XLSX) the model cannot parse the raw bytes — this tool does the
   decoding and returns plain text so the model can answer questions
   about the content immediately, no GCS required.

2. save_attachment_for_indexing(file_base64, filename, folder)
   ────────────────────────────────────────────────────────────
   Uploads any attached file to the GCS bucket under an "uploads/"
   prefix and triggers the full ingestion pipeline so the document
   becomes searchable in future queries.  Supported formats: same as
   the ingestion pipeline (PDF, DOCX, PPTX, XLSX, TXT, MD, CSV).

Note: no 'from __future__ import annotations' — ADK runtime introspection.
"""

import base64
import io
import logging
import os
import tempfile
from datetime import datetime

from google.adk.tools import ToolContext

logger = logging.getLogger(__name__)

_UPLOAD_BUCKET = "fre-cognitive-search-docs"
_UPLOAD_ROOT   = "User_Uploads"


# ── helpers ───────────────────────────────────────────────────────────────────

def _ext(filename: str) -> str:
    return os.path.splitext(filename.lower())[1]


def _extract_docx(data: bytes) -> str:
    import docx  # python-docx
    doc = docx.Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # also grab table cells
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n".join(paragraphs)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    lines = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(v) for v in row if v is not None)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)


# ── public tools ──────────────────────────────────────────────────────────────

def extract_office_document_text(file_base64: str, filename: str) -> dict:
    """
    Extract plain text from an Office document (DOCX, PPTX, or XLSX) that
    the user has attached to the chat.

    Gemini can read PDF and image attachments natively, but cannot parse
    binary Office formats.  Call this tool when the user attaches a
    .docx, .pptx, or .xlsx file and asks a question about its contents.

    Parameters
    ----------
    file_base64 : Base64-encoded file content (the raw bytes of the attachment).
    filename    : Original filename including extension, e.g. "report.docx".
                  Used to determine the file format.

    Returns
    -------
    Dict with keys:
      filename      (str)  : Echo of the input filename.
      format        (str)  : Detected format ("docx" | "pptx" | "xlsx" | "unsupported").
      text          (str)  : Extracted plain text (empty string on failure).
      char_count    (int)  : Number of characters extracted.
      status        (str)  : "ok" | "failed" | "unsupported_format".
      message       (str)  : Human-readable result summary.
    """
    result: dict = {
        "filename":   filename,
        "format":     "unknown",
        "text":       "",
        "char_count": 0,
        "status":     "failed",
        "message":    "",
    }

    try:
        data = base64.b64decode(file_base64)
    except Exception as exc:
        result["message"] = f"Base64 decode failed: {exc}"
        return result

    ext = _ext(filename)
    result["format"] = ext.lstrip(".")

    try:
        if ext == ".docx":
            text = _extract_docx(data)
        elif ext == ".pptx":
            text = _extract_pptx(data)
        elif ext in (".xlsx", ".xls"):
            text = _extract_xlsx(data)
        else:
            result["status"] = "unsupported_format"
            result["message"] = (
                f"Format '{ext}' is not an Office binary format. "
                "PDF and image files are read natively by Gemini — "
                "no extraction tool needed."
            )
            return result

        result["text"] = text
        result["char_count"] = len(text)
        result["status"] = "ok"
        result["message"] = (
            f"Extracted {len(text):,} characters from {filename}. "
            "The full text is in the 'text' field — use it to answer the user's question."
        )
        logger.info(
            "extract_office_document_text: %s → %d chars", filename, len(text)
        )

    except Exception as exc:
        result["message"] = f"Extraction failed: {exc}"
        logger.error("extract_office_document_text error for %s: %s", filename, exc)

    return result


def save_attachment_for_indexing(
    file_base64: str,
    filename: str,
    folder: str = "",
) -> dict:
    """
    Save a user-attached file to GCS and index it so it becomes searchable
    in future queries.

    Files are stored under:
      gs://fre-cognitive-search-docs/User_Uploads/YYYY-MM/<folder>/<filename>

    Where:
      • YYYY-MM  is automatically set to the current year-month.
      • <folder> is an optional sub-folder the user can specify to group
                 related uploads (e.g. "Project_Aurora", "Customer_Docs").
                 Leave empty to place files directly in the month folder.

    Use this tool when the user attaches a document and says something like:
      - "save this to the knowledge base"
      - "index this file"
      - "add this document to the search index"
      - "make this searchable"

    After a successful upload and index, the document will appear in all
    future hybrid searches (both Elasticsearch and Vertex AI).

    Parameters
    ----------
    file_base64 : Base64-encoded file content.
    filename    : Original filename including extension, e.g. "report.pdf".
    folder      : Optional sub-folder name inside the month folder.
                  Example: "Project_Aurora" → stored at
                  User_Uploads/2026-04/Project_Aurora/report.pdf
                  Leave empty to store directly in User_Uploads/2026-04/.

    Returns
    -------
    Dict with keys:
      gcs_uri       (str)  : gs:// URI where the file was saved.
      filename      (str)  : Echo of filename.
      size_kb       (float): File size in kilobytes.
      indexed       (bool) : Whether ingestion was triggered successfully.
      status        (str)  : "saved_and_indexed" | "saved_only" | "failed".
      message       (str)  : Human-readable result summary.
    """
    result: dict = {
        "gcs_uri":  "",
        "filename": filename,
        "size_kb":  0.0,
        "indexed":  False,
        "status":   "failed",
        "message":  "",
    }

    # ── decode ────────────────────────────────────────────────────────────────
    try:
        data = base64.b64decode(file_base64)
    except Exception as exc:
        result["message"] = f"Base64 decode failed: {exc}"
        return result

    result["size_kb"] = round(len(data) / 1024, 1)

    # ── upload to GCS ─────────────────────────────────────────────────────────
    try:
        from storage.gcs import upload_bytes

        # Build path: User_Uploads/YYYY-MM[/folder]/filename
        month_prefix = datetime.now().strftime("%Y-%m")
        parts = [_UPLOAD_ROOT, month_prefix]
        if folder.strip():
            parts.append(folder.strip().replace("/", "_"))
        parts.append(filename)
        blob_path = "/".join(parts)

        ext = _ext(filename)
        mime_map = {
            ".pdf":  "application/pdf",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".txt":  "text/plain",
            ".md":   "text/markdown",
            ".csv":  "text/csv",
            ".png":  "image/png",
            ".jpg":  "image/jpeg",
            ".jpeg": "image/jpeg",
        }
        content_type = mime_map.get(ext, "application/octet-stream")

        gcs_uri = upload_bytes(data, _UPLOAD_BUCKET, blob_path, content_type)
        result["gcs_uri"] = gcs_uri
        logger.info("save_attachment_for_indexing: uploaded %s → %s", filename, gcs_uri)

    except Exception as exc:
        result["message"] = f"GCS upload failed: {exc}"
        logger.error("save_attachment_for_indexing upload error: %s", exc)
        return result

    # ── trigger ingestion ─────────────────────────────────────────────────────
    try:
        from tools.ingestion_tools import trigger_document_ingestion
        ingestion_result = trigger_document_ingestion(gcs_uri)
        result["indexed"] = ingestion_result.get("status") not in ("failed", "error")
        result["status"] = "saved_and_indexed" if result["indexed"] else "saved_only"
        result["message"] = (
            f"Uploaded {filename} ({result['size_kb']} KB) to {gcs_uri}. "
            + (
                f"Ingestion triggered: {ingestion_result.get('message', '')}"
                if result["indexed"]
                else f"Upload succeeded but indexing failed: {ingestion_result.get('message', '')}"
            )
        )
    except Exception as exc:
        result["status"] = "saved_only"
        result["message"] = (
            f"Uploaded {filename} ({result['size_kb']} KB) to {gcs_uri} "
            f"but could not trigger ingestion: {exc}"
        )
        logger.error("save_attachment_for_indexing ingestion error: %s", exc)

    return result


def load_attachment_to_session(
    file_base64: str,
    filename: str,
    tool_context: ToolContext,  # ADK injects this automatically
) -> str:
    """
    Load an uploaded document into the current session so it can be
    referenced in follow-up questions WITHOUT re-uploading the file.

    After calling this tool the document content is stored in session state
    and remains available for all subsequent questions in this conversation.

    Use this tool when the user:
      - Attaches a file AND asks a question about it
      - Wants to discuss a document across multiple messages ("ask me anything
        about this PDF", "I'll send a file, analyse it")

    Supported formats: PDF, DOCX, PPTX, XLSX, TXT, MD, CSV

    For PDF/image files Gemini can also read the attachment natively in the
    same message — but calling this tool additionally stores the text for
    follow-up questions.

    Parameters
    ----------
    file_base64  : Base64-encoded file content (provided by the ADK UI).
    filename     : Original filename including extension, e.g. "report.pdf".
    tool_context : Injected by ADK — do NOT pass this from the LLM.

    Returns
    -------
    A text string confirming the document was loaded, including the first
    2 000 characters of the extracted content so you can answer the user's
    first question immediately.
    """
    try:
        data = base64.b64decode(file_base64)
    except Exception as exc:
        return f"ERROR: Base64 decode failed for {filename}: {exc}"

    ext = _ext(filename)
    size_kb = round(len(data) / 1024, 1)

    # ── extract text ──────────────────────────────────────────────────────────
    text = ""
    try:
        if ext == ".pdf":
            text = _extract_pdf_text(data)
        elif ext == ".docx":
            text = _extract_docx(data)
        elif ext == ".pptx":
            text = _extract_pptx(data)
        elif ext in (".xlsx", ".xls"):
            text = _extract_xlsx(data)
        elif ext in (".txt", ".md", ".csv"):
            text = data.decode("utf-8", errors="replace")
        else:
            return (
                f"Format '{ext}' cannot be extracted as text. "
                "PDF and images are read natively by Gemini in the same message. "
                "For follow-up questions, re-attach the file."
            )
    except Exception as exc:
        return f"ERROR extracting text from {filename}: {exc}"

    if not text.strip():
        return f"WARNING: No text could be extracted from {filename} ({size_kb} KB)."

    # ── store in session state ────────────────────────────────────────────────
    # ADK merges tool_context.state into session state after each tool call.
    # The agent sees all stored documents in subsequent turns via session state.
    existing: dict = tool_context.state.get("session_documents", {})
    existing[filename] = text
    tool_context.state["session_documents"] = existing

    # also keep a human-readable list for the agent's context
    doc_list = list(existing.keys())
    tool_context.state["session_document_list"] = doc_list

    logger.info(
        "load_attachment_to_session: stored '%s' (%d chars) in session state",
        filename, len(text),
    )

    preview = text[:2000].strip()
    return (
        f"✅ **{filename}** ({size_kb} KB, {len(text):,} chars) loaded into this session.\n\n"
        f"You can ask follow-up questions about it at any time — it is stored in session memory.\n\n"
        f"**Document preview (first 2 000 chars):**\n\n{preview}"
        + ("\n\n*(document continues...)*" if len(text) > 2000 else "")
    )


def _extract_pdf_text(data: bytes) -> str:
    """Extract plain text from PDF bytes using pypdf (no Vision, fast)."""
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(data))
        pages = []
        for i, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                pages.append(f"[Page {i}]\n{page_text}")
        return "\n\n".join(pages)
    except ImportError:
        pass

    # fallback: PyMuPDF
    try:
        import fitz
        doc = fitz.open(stream=data, filetype="pdf")
        pages = []
        for i, page in enumerate(doc, start=1):
            t = page.get_text()
            if t.strip():
                pages.append(f"[Page {i}]\n{t}")
        doc.close()
        return "\n\n".join(pages)
    except Exception as exc:
        raise RuntimeError(f"PDF text extraction failed: {exc}") from exc
