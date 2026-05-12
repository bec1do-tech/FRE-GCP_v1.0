"""
FRE GCP v1.0 — Multimodal Document Extractor
==============================================
Extracts text content from all supported file types, including:
  • Plain text / Markdown / CSV         → direct decode
  • PDF                                  → pypdf (text) + Gemini Vision (images/charts)
  • DOCX                                 → python-docx (text + embedded images)
  • PPTX                                 → python-pptx (slide text + Gemini Vision per slide)
  • XLSX                                 → openpyxl (table data as structured text)

For binary-heavy files (PDF/PPTX/DOCX) each embedded image is sent to the
Gemini Vision model to produce a natural-language description, which is then
appended to the extracted text.  This means charts, diagrams, and infographics
become part of the searchable corpus.

All functions accept raw bytes so they work identically with GCS downloads
and local file reads.
"""

from __future__ import annotations

import io
import logging
import zipfile
import re
from pathlib import PurePosixPath
from typing import NamedTuple

import config

logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────────────────────────────────────
# Return type
# ─────────────────────────────────────────────────────────────────────────────

class ExtractionResult(NamedTuple):
    text: str               # full extracted text (including image descriptions)
    metadata: dict          # file-level metadata (author, page_count, etc.)
    image_count: int        # how many images were processed via Vision


# ─────────────────────────────────────────────────────────────────────────────
# Gemini Vision helper
# ─────────────────────────────────────────────────────────────────────────────

def _genai_client():
    """Return a google.genai Client pointed at Vertex AI, reused per process."""
    from google import genai  # type: ignore[import-untyped]

    return genai.Client(
        vertexai=True,
        project=config.GCP_PROJECT,
        location=config.GCP_REGION,
    )


def _describe_image(image_bytes: bytes, mime_type: str = "image/png") -> str:
    """
    Send one image to Gemini Vision (via Vertex AI) and return a description.
    Returns "" if Gemini is unavailable or the image is too small.
    """
    import concurrent.futures

    if len(image_bytes) < 1024:   # skip trivially small images (bullets, icons)
        return ""

    def _call() -> str:
        from google.genai import types  # type: ignore[import-untyped]

        client = _genai_client()
        response = client.models.generate_content(
            model=config.GEMINI_VISION_MODEL,
            contents=[
                "Describe the content of this image in detail. "
                "If it is a chart or graph, summarise the key data points and trends. "
                "If it is a table, transcribe the data. "
                "If it is a diagram, describe its structure and labels. "
                "Be concise but complete.",
                types.Part.from_bytes(data=image_bytes, mime_type=mime_type),
            ],
        )
        return (response.text or "").strip()

    timeout_s = int(os.environ.get("VISION_TIMEOUT_S", "60"))
    try:
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
            future = pool.submit(_call)
            return future.result(timeout=timeout_s)
    except concurrent.futures.TimeoutError:
        logger.warning("Gemini Vision timed out after %ds — skipping image.", timeout_s)
        return ""
    except Exception as exc:
        logger.debug("Gemini Vision unavailable: %s", exc)
        return ""


def _ocr_pdf_with_gemini(data: bytes) -> str:
    """
    OCR fallback for scanned/image-only PDFs.
    Sends the entire PDF bytes to Gemini via Vertex AI in a single call and
    asks it to transcribe all visible text.
    Returns "" on failure, timeout, or empty response.
    Skips files > OCR_MAX_PDF_MB to avoid inline-data limits.
    """
    import concurrent.futures

    ocr_limit_mb = getattr(config, "OCR_MAX_PDF_MB", 18)  # Gemini inline limit ~20 MB
    ocr_timeout_s = int(os.environ.get("OCR_TIMEOUT_S", "120"))  # max seconds to wait
    size_mb = len(data) / (1024 * 1024)
    if size_mb > ocr_limit_mb:
        logger.warning(
            "PDF too large for Gemini OCR (%.1f MB > %d MB limit) — skipping OCR.",
            size_mb, ocr_limit_mb,
        )
        return ""

    def _call() -> str:
        from google.genai import types  # type: ignore[import-untyped]

        client = _genai_client()
        response = client.models.generate_content(
            model=config.GEMINI_VISION_MODEL,
            contents=[
                "This PDF is a scanned document with no embedded text layer. "
                "Transcribe ALL visible text exactly as it appears, page by page. "
                "Begin each page with a [Page N] marker. "
                "For tables, transcribe the content row by row. "
                "For pages with no readable text write [No text on page N]. "
                "Output only the transcribed text — no commentary.",
                types.Part.from_bytes(data=data, mime_type="application/pdf"),
            ],
        )
        return (response.text or "").strip()

    try:
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
            future = pool.submit(_call)
            text = future.result(timeout=ocr_timeout_s)
        logger.info("Gemini OCR produced %d characters from PDF.", len(text))
        return text
    except concurrent.futures.TimeoutError:
        logger.warning("Gemini PDF OCR timed out after %ds.", ocr_timeout_s)
        return ""
    except Exception as exc:
        logger.warning("Gemini PDF OCR failed: %s", exc)
        return ""


# ─────────────────────────────────────────────────────────────────────────────
# Plain text / Markdown / CSV
# ─────────────────────────────────────────────────────────────────────────────

def _extract_text_bytes(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except (UnicodeDecodeError, AttributeError):
            continue
    return data.decode("utf-8", errors="replace")


# ─────────────────────────────────────────────────────────────────────────────
# PDF
# ─────────────────────────────────────────────────────────────────────────────

def _extract_pdf(data: bytes) -> ExtractionResult:
    try:
        import pypdf  # type: ignore[import-untyped]
        import pypdf.errors  # type: ignore[import-untyped]
    except ImportError:
        logger.warning("pypdf not installed — PDF text extraction disabled.")
        return ExtractionResult("", {}, 0)

    text_parts: list[str] = []
    image_descriptions: list[str] = []
    metadata: dict = {}

    try:
        reader = pypdf.PdfReader(io.BytesIO(data))
        metadata = {
            "page_count": len(reader.pages),
            "author":     (reader.metadata or {}).get("/Author", ""),
            "title":      (reader.metadata or {}).get("/Title", ""),
            "creator":    (reader.metadata or {}).get("/Creator", ""),
        }

        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"[Page {page_num + 1}]\n{page_text}")

            # Extract embedded images and describe with Gemini Vision
            for img_obj in page.images:
                try:
                    desc = _describe_image(img_obj.data, "image/png")
                    if desc:
                        image_descriptions.append(
                            f"[Image on page {page_num + 1}: {desc}]"
                        )
                except Exception:
                    pass

    except Exception as exc:
        logger.error("PDF extraction error: %s", exc)

    full_text = "\n\n".join(text_parts)
    if image_descriptions:
        full_text += "\n\n--- Visual Content ---\n" + "\n".join(image_descriptions)

    # OCR fallback: if pypdf found no text, the PDF is likely a scanned image.
    if not full_text.strip():
        logger.info("No text layer detected — attempting Gemini OCR fallback.")
        ocr_text = _ocr_pdf_with_gemini(data)
        if ocr_text:
            full_text = ocr_text
            metadata["ocr"] = True

    return ExtractionResult(full_text, metadata, len(image_descriptions))


# ─────────────────────────────────────────────────────────────────────────────
# DOCX
# ─────────────────────────────────────────────────────────────────────────────

def _extract_docx(data: bytes) -> ExtractionResult:
    try:
        from docx import Document  # type: ignore[import-untyped]
        from docx.oxml.ns import qn  # type: ignore[import-untyped]
    except ImportError:
        logger.warning("python-docx not installed — DOCX extraction disabled.")
        return ExtractionResult("", {}, 0)

    text_parts: list[str] = []
    image_descriptions: list[str] = []
    metadata: dict = {}

    try:
        doc = Document(io.BytesIO(data))
        core = doc.core_properties
        metadata = {
            "author":   core.author or "",
            "title":    core.title or "",
            "subject":  core.subject or "",
            "created":  core.created.isoformat() if core.created else "",
        }

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Extract inline images
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    img_bytes = rel.target_part.blob
                    desc = _describe_image(img_bytes)
                    if desc:
                        image_descriptions.append(f"[Document image: {desc}]")
                except Exception:
                    pass

    except Exception as exc:
        logger.error("DOCX extraction error: %s", exc)

    full_text = "\n".join(text_parts)
    if image_descriptions:
        full_text += "\n\n--- Visual Content ---\n" + "\n".join(image_descriptions)

    return ExtractionResult(full_text, metadata, len(image_descriptions))


# ─────────────────────────────────────────────────────────────────────────────
# PPTX
# ─────────────────────────────────────────────────────────────────────────────

def _extract_pptx(data: bytes) -> ExtractionResult:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.util import Pt  # type: ignore[import-untyped]
    except ImportError:
        logger.warning("python-pptx not installed — PPTX extraction disabled.")
        return ExtractionResult("", {}, 0)

    text_parts: list[str] = []
    image_descriptions: list[str] = []
    metadata: dict = {}

    try:
        prs = Presentation(io.BytesIO(data))
        metadata = {"slide_count": len(prs.slides)}

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs if run.text.strip())
                        if line:
                            slide_texts.append(line)

                # Images on slide
                if shape.shape_type == 13:   # MSO_SHAPE_TYPE.PICTURE
                    try:
                        img_bytes = shape.image.blob
                        # Render the whole slide as PNG for richer context
                        desc = _describe_image(img_bytes)
                        if desc:
                            image_descriptions.append(
                                f"[Image on slide {slide_num}: {desc}]"
                            )
                    except Exception:
                        pass

            if slide_texts:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

    except Exception as exc:
        logger.error("PPTX extraction error: %s", exc)

    full_text = "\n\n".join(text_parts)
    if image_descriptions:
        full_text += "\n\n--- Visual Content ---\n" + "\n".join(image_descriptions)

    return ExtractionResult(full_text, metadata, len(image_descriptions))


# ─────────────────────────────────────────────────────────────────────────────
# XLSX
# ─────────────────────────────────────────────────────────────────────────────

def _extract_xlsx(data: bytes) -> ExtractionResult:
    try:
        import openpyxl  # type: ignore[import-untyped]
    except ImportError:
        logger.warning("openpyxl not installed — XLSX extraction disabled.")
        return ExtractionResult("", {}, 0)

    text_parts: list[str] = []
    metadata: dict = {}

    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        metadata = {
            "sheet_names": wb.sheetnames,
            "sheet_count": len(wb.sheetnames),
        }

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

        wb.close()
    except Exception as exc:
        logger.error("XLSX extraction error: %s", exc)

    return ExtractionResult("\n\n".join(text_parts), metadata, 0)


# ─────────────────────────────────────────────────────────────────────────────
# Dispatcher
# ─────────────────────────────────────────────────────────────────────────────

_EXTRACTORS = {
    ".txt":  lambda d: ExtractionResult(_extract_text_bytes(d), {}, 0),
    ".md":   lambda d: ExtractionResult(_extract_text_bytes(d), {}, 0),
    ".csv":  lambda d: ExtractionResult(_extract_text_bytes(d), {}, 0),
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
}

SUPPORTED_EXTENSIONS: tuple[str, ...] = tuple(_EXTRACTORS.keys())


def extract(data: bytes, filename: str) -> ExtractionResult:
    """
    Extract text (and visual descriptions) from raw file bytes.

    Parameters
    ----------
    data     : Raw file bytes (e.g. from GCS download).
    filename : Original filename — used only to determine the file type.

    Returns
    -------
    ExtractionResult(text, metadata, image_count)
    """
    ext = PurePosixPath(filename).suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        logger.warning("Unsupported file type: %s", filename)
        return ExtractionResult("", {}, 0)

    mb = len(data) / (1024 * 1024)
    if mb > config.MAX_FILE_MB:
        logger.warning("File %s (%.1f MB) exceeds MAX_FILE_MB=%d — skipping.", filename, mb, config.MAX_FILE_MB)
        return ExtractionResult("", {}, 0)

    return extractor(data)
